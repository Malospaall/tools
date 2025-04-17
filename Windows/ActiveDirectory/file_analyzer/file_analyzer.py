import os
import re
import argparse
import platform
from docx import Document
from xml.etree import ElementTree as ET
import pdfplumber
from pdf2image import convert_from_path
import pytesseract
import pptx
import xlrd
import openpyxl
import json
import jsbeautifier
from striprtf.striprtf import rtf_to_text
import email
import extract_msg
import chardet
import shutil
import tempfile
import zipfile
import rarfile
import py7zr
import concurrent.futures

# Поиск типа файла
def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    file_readers = {
        '.doc': read_doc, '.docx': read_docx, '.dotx': read_dotx,
        '.pdf': read_pdf,
        '.pptx': read_pptx,
        '.xls': read_xls,
        '.xlsx': read_xlsx,
        '.jpg': read_image, '.jpeg': read_image, '.png': read_image, '.tiff': read_image, '.bmp': read_image,
        '.txt': read_txt, '.md': read_txt, '.log': read_txt, '.config': read_txt, '.conf': read_txt, '.cfg': read_txt, '.ini': read_txt, '.ps1': read_txt,
        '.sh': read_txt, '.csv': read_txt, '.reg': read_txt, '.py': read_txt, '.cs': read_txt, '.c': read_txt, '.asp': read_txt,
        '.aspx': read_txt, '.ashx': read_txt, '.php': read_txt, '.xml': read_txt, '.html': read_txt, '.htm': read_txt, '.css': read_txt, '': read_txt, '.lnk': read_txt,
        '.js': read_js,
        '.json': read_json,
        '.rtf': read_rtf,
        '.eml': read_eml,
        '.msg': read_msg,
        '.exe': read_execute, '.dll': read_execute,
    }
    reader = file_readers.get(ext)
    if reader:
        return reader(file_path)
    return f"Unsupported file type: {ext}"

# Обработка файла
def read_doc(file_path):
    system = platform.system()

    if system == "Windows":
        try:
            import comtypes.client
            word = comtypes.client.CreateObject("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(file_path)
            text = doc.Content.Text
            doc.Close(False)
            word.Quit()
            return text
        except Exception as e:
            return f"Error reading DOC file on Windows: {e}"

    elif system == "Linux":
        try:
            import subprocess
            text = subprocess.check_output(['catdoc', file_path], stderr=subprocess.PIPE).decode('utf-8')
            return text
        except Exception as e:
            return f"Error reading DOC file on Linux: {e}"

def read_docx(file_path):
    doc = Document(file_path)
    content = []
    for para in doc.paragraphs:
        content.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text)
            content.append('\t'.join(row_data))
    return '\n'.join(content)

def read_dotx(file_path):
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            with zip_ref.open('word/document.xml') as document_file:
                tree = ET.parse(document_file)
                root = tree.getroot()
                namespaces = {
                    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
                }
                text_content = []
                for paragraph in root.findall('.//w:p', namespaces):
                    text = ''.join(node.text for node in paragraph.findall('.//w:t', namespaces) if node.text)
                    text_content.append(text)
                return '\n'.join(text_content)
    except Exception as e:
        return f"Error reading DOTX file: {e}"

def read_pdf(file_path):
    with pdfplumber.open(file_path) as pdf:
        text = '\n'.join([page.extract_text() for page in pdf.pages if page.extract_text()])
        if not text.strip():
            text = read_pdf_with_ocr(file_path)
        return text

def read_pdf_with_ocr(file_path):
    try:
        images = convert_from_path(file_path)
        return '\n'.join([pytesseract.image_to_string(img, lang='rus+eng') for img in images])
    except Exception as e:
        return f"Error extracting text with OCR: {e}"

def read_pptx(file_path):
    presentation = pptx.Presentation(file_path)
    return '\n'.join([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, 'text')])

def read_xls(file_path):
    wb = xlrd.open_workbook(file_path)
    sheet = wb.sheet_by_index(0)
    return '\n'.join(['\t'.join(map(str, sheet.row_values(row))) for row in range(sheet.nrows)])

def read_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    sheet = wb.active
    return '\n'.join(['\t'.join([str(cell.value) for cell in row]) for row in sheet.iter_rows()])

def read_image(file_path):
    try:
        return pytesseract.image_to_string(file_path, lang='rus+eng')
    except Exception as e:
        return f"Error extracting text from image: {e}"

def read_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"Error reading text file: {e}"

def read_json(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, indent=4, ensure_ascii=False)
    except Exception as e:
        return f"Error reading JSON file: {e}"

def read_js(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            js_code = f.read()

        beautifier = jsbeautifier.Beautifier(js_code)
        formatted_js = jsbeautifier.beautify(js_code)

        return formatted_js
    except Exception as e:
        return f"Error reading or formatting JS file: {e}"

def read_rtf(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return rtf_to_text(f.read())
    except Exception as e:
        return f"Error reading RTF file: {e}"

def read_eml(file_path):
    try:
        with open(file_path, "rb") as f:
            raw_data = f.read()
        encoding = chardet.detect(raw_data).get("encoding", "utf-8")
        msg = email.message_from_bytes(raw_data)
        text = "\n".join(part.get_payload(decode=True).decode(encoding, errors="ignore") for part in msg.walk() if part.get_content_type() == "text/plain")
        return text if text else "No text content"
    except Exception as e:
        return f"Error reading EML file: {e}"

def read_msg(file_path):
    try:
        msg = extract_msg.Message(file_path)
        body = msg.body
        encoding = chardet.detect(body.encode()).get("encoding", "utf-8")
        return f"Отправитель: {msg.sender}\nДата: {msg.date}\nТема: {msg.subject}\nТело письма:\n\n{body.encode().decode(encoding, errors='ignore')}"
    except Exception as e:
        return f"Error reading MSG file: {e}"

def read_execute(file_path):
    try:
        with open(file_path, 'rb') as f:
            content = f.read()
            return '\n'.join(
                ''.join(chr(byte) for byte in content[i:i+80] if 32 <= byte <= 126)
                for i in range(0, len(content), 80)
            )
    except Exception as e:
        return f"Error reading executable file: {e}"

# Обработка архива
def extract_zip(file_path, extract_to):
    return extract_archive(file_path, extract_to, zipfile.ZipFile)

def extract_rar(file_path, extract_to):
    return extract_archive(file_path, extract_to, rarfile.RarFile)

def extract_7z(file_path, extract_to):
    return extract_archive(file_path, extract_to, py7zr.SevenZipFile)

def extract_archive(file_path, extract_to, archive_type):
    try:
        with archive_type(file_path, 'r') as archive_ref:
            archive_ref.extractall(extract_to)
        return extract_to
    except Exception as e:
        return f"Error extracting archive: {e}"

# Формирование HTML страницы
def generate_html(output_file):
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write("""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <title>File Extraction</title>
            <style>
                body { font-family: Arial, sans-serif; margin: 20px; }
                h1 { color: #2c3e50; }
                input { width: 300px; padding: 5px; margin-bottom: 20px; }
                pre { background: #f4f4f4; padding: 10px; border-radius: 5px; overflow-x: auto; }
                .file-container { border-bottom: 1px solid #ddd; padding: 10px 0; }
                .file-name { font-weight: bold; color: #007bff; cursor: pointer; }
                .file-block { margin-bottom: 20px; border-bottom: 1px solid #ccc; padding-bottom: 10px; }
            </style>
            <script>
                function searchFiles() {
                    let query = document.getElementById('search').value.toLowerCase();
                    let files = document.getElementsByClassName('file-container');

                    for (let file of files) {
                        let name = file.getElementsByClassName('file-name')[0].innerText.toLowerCase();
                        let content = file.getElementsByClassName('file-content')[0].innerText.toLowerCase();

                        file.style.display = (name.includes(query) || content.includes(query)) ? 'block' : 'none';
                    }
                }

                function toggleContent(id) {
                    let content = document.getElementById(id);
                    content.style.display = content.style.display === 'none' ? 'block' : 'none';
                }
            </script>
        </head>
        <body>
            <h1>File Extraction</h1>
            <input type="text" id="search" onkeyup="searchFiles()" placeholder="search file...">
        """)


# Добавление файлов в HTML
def append_file_to_html(file_id, file_path, content, output_file):
    escaped_content = content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;").replace("'", "&#39;")
    with open(output_file, 'a', encoding='utf-8') as f:
        f.write(f"""
        <div class="file-container">
            <div class="file-name" onclick="toggleContent('{file_id}')">{file_path}</div>
            <pre class="file-content" id="{file_id}" style="display:none;">{escaped_content}</pre>
        </div>
        """)

def finalize_html(output_file):
    with open(output_file, 'a', encoding='utf-8') as f:
        f.write("""</body></html>""")
    print(f"\n[+] HTML file: {output_file}")

# Парсинг размера файла
def parse_size(size_str):
    match = re.match(r"(\d+)(B|KB|MB|GB)", size_str, re.IGNORECASE)
    if not match:
        raise ValueError("Invalid size format. Use formats like 10KB, 10MB, 1GB, etc.")
    size, unit = match.groups()
    size = int(size)
    unit_multipliers = {"B": 1, "KB": 1024, "MB": 1024**2, "GB": 1024**3}
    return size * unit_multipliers[unit.upper()]

# Функция для обработки файла
def process_file(file_id, file_path):
    try:
        content = read_file(file_path)
        return file_id, file_path, content
    except Exception as e:
        return file_id, file_path, f"Error: {e}"

# Сбор файлов в директории
def process_folder(folder_path, exclude_dirs, exclude_ext, output_file):
    folder_path = os.path.abspath(folder_path)
    temp_dir = tempfile.mkdtemp()
    all_files = []

    if not os.path.exists(folder_path) or not os.path.isdir(folder_path):
        print(f"Error: Directory '{folder_path}' does not exist or is not accessible.")
        exit(1)

    generate_html(output_file)

    def collect_files(extracted_path):
        # Рекурсивный сбор файлов из архива
        for root, _, files in os.walk(extracted_path):
            for file in files:
                file_path = os.path.join(root, file)
                if max_size and os.path.getsize(file_path) > max_size:
                    continue
                all_files.append(file_path)

    # Обрабатываем файлы и архивы в папке
    for root, dirs, files in os.walk(folder_path):
        # Исключение директорий
        dirs[:] = [d for d in dirs if d not in exclude_dirs]
        # Глубина вложенности
        depth = root[len(folder_path):].count(os.sep)
        if args.depth is not None and depth >= args.depth:
            dirs.clear()

        for file in files:
            file_path = os.path.join(root, file)
            ext = os.path.splitext(file)[1].lower()

            if exclude_ext and ext in exclude_ext:
                continue

            if max_size and os.path.getsize(file_path) > max_size:
                continue

            if ext in ['.zip', '.rar', '.7z']:
                extracted_path = globals().get(f"extract_{ext[1:]}")(file_path, temp_dir)
                if os.path.isdir(extracted_path):
                    # Рекурсивно собираем файлы из архива
                    collect_files(extracted_path)
            else:
                all_files.append(file_path)

    # Обработка файлов
    with concurrent.futures.ProcessPoolExecutor() as executor:
        futures = {executor.submit(process_file, idx, file_path): (idx, file_path) for idx, file_path in enumerate(all_files)}
        for i, future in enumerate(concurrent.futures.as_completed(futures), start=1):
            file_id, file_path, content = future.result()
            append_file_to_html(file_id, file_path, content, output_file)
            print(f"\rProcessing: {i} / {len(all_files)}", end='', flush=True)

    finalize_html(output_file)
    shutil.rmtree(temp_dir)
    print("[+] Processing complete.")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Recursive file reader and extractor")
    parser.add_argument("-d", dest="directory", required=True, help="Path to the folder containing files")
    parser.add_argument("-e", dest="exclude", help="List of directories to exclude")
    parser.add_argument("-s", dest="size", help="Maximum file size. Example: 10KB, 10MB, 1GB (default: no limit)", default=None)
    parser.add_argument("-x", dest="exclude_ext", help="List of file extensions to exclude", default="")
    parser.add_argument("-depth", dest="depth", type=int, help="Maximum folder depth for recursive scanning (default: no limit)")
    parser.add_argument("-o", dest="output", help="Output html file name (default: output.html)", default="output.html")
    args = parser.parse_args()

    exclude_dirs = args.exclude.split(',') if args.exclude else []
    max_size = parse_size(args.size) if args.size else None
    exclude_ext = set(ext.strip().lower() for ext in args.exclude_ext.split(',') if ext.strip())

    process_folder(args.directory, exclude_dirs, exclude_ext, args.output)